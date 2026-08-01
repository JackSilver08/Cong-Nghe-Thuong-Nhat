from pathlib import Path
from io import BytesIO
from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

ROOT = Path(__file__).resolve().parents[1]
OUT = ROOT / "docs" / "Thuyet-trinh-Cong-Nghe-Thuong-Nhat-5-slide.pptx"

NAVY = RGBColor(9, 17, 39)
BLUE = RGBColor(50, 69, 232)
CYAN = RGBColor(35, 211, 238)
WHITE = RGBColor(255, 255, 255)
INK = RGBColor(17, 24, 39)
MUTED = RGBColor(91, 101, 120)
PALE = RGBColor(244, 247, 252)
LINE = RGBColor(218, 225, 236)

prs = Presentation()
prs.slide_width = Inches(13.333333)
prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]


def rect(slide, x, y, w, h, color, radius=False, line=None, transparency=0):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE,
        Inches(x), Inches(y), Inches(w), Inches(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.fill.transparency = transparency
    shape.line.color.rgb = line or color
    return shape


def text(slide, value, x, y, w, h, size=20, color=INK, bold=False,
         font="Arial", align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP, margin=0):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.margin_left = frame.margin_right = Inches(margin)
    frame.margin_top = frame.margin_bottom = Inches(margin)
    frame.vertical_anchor = valign
    p = frame.paragraphs[0]
    p.text = value
    p.alignment = align
    p.font.name = font
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.space_after = Pt(0)
    return box


def add_image_crop(slide, path, x, y, w, h):
    path = Path(path)
    with Image.open(path) as im:
        iw, ih = im.size
        image_source = str(path)
        if im.format == "WEBP":
            image_source = BytesIO()
            im.convert("RGB").save(image_source, format="PNG")
            image_source.seek(0)
    target = w / h
    source = iw / ih
    picture = slide.shapes.add_picture(image_source, Inches(x), Inches(y), Inches(w), Inches(h))
    if source > target:
        shown = ih * target
        crop = (iw - shown) / 2 / iw
        picture.crop_left = crop
        picture.crop_right = crop
    else:
        shown = iw / target
        crop = (ih - shown) / 2 / ih
        picture.crop_top = crop
        picture.crop_bottom = crop
    return picture


def footer(slide, index, dark=False):
    color = RGBColor(171, 184, 213) if dark else MUTED
    text(slide, "CÔNG NGHỆ THƯỜNG NHẬT", .55, 7.12, 3.1, .2, 8, color, True)
    text(slide, f"0{index} / 05", 11.95, 7.08, .8, .24, 8, color, True, align=PP_ALIGN.RIGHT)


# Slide 1 — Cover
s = prs.slides.add_slide(blank)
add_image_crop(s, ROOT / "public/assets/login-tech-news-hero.webp", 0, 0, 13.333, 7.5)
rect(s, 0, 0, 13.333, 7.5, NAVY, transparency=18)
rect(s, 0, 0, 7.8, 7.5, NAVY, transparency=5)
text(s, "CÔNG NGHỆ\nTHƯỜNG NHẬT", .72, .62, 3.8, .85, 21, WHITE, True)
text(s, "NỀN TẢNG TIN TỨC CÔNG NGHỆ", .75, 2.1, 4.4, .3, 11, CYAN, True)
text(s, "Công nghệ dễ hiểu,\nhữu ích mỗi ngày.", .72, 2.5, 7.2, 1.6, 34, WHITE, True)
text(s, "Giới thiệu website & trải nghiệm số", .75, 4.38, 5.6, .4, 17, RGBColor(218, 226, 244))
rect(s, .75, 5.2, 2.28, .48, BLUE, radius=True)
text(s, "BÀI THUYẾT TRÌNH  •  2026", .88, 5.34, 2.03, .2, 9, WHITE, True)
footer(s, 1, True)

# Slide 2 — Purpose
s = prs.slides.add_slide(blank)
rect(s, 0, 0, 13.333, 7.5, PALE)
text(s, "01 / SỨ MỆNH", .62, .5, 2.2, .25, 10, BLUE, True)
text(s, "Biến tin công nghệ phức tạp\nthành nội dung dễ tiếp cận.", .62, .93, 6.15, 1.25, 27, INK, True)
text(s, "Một điểm đến cho người Việt muốn hiểu công nghệ\nnhanh, đúng trọng tâm và có thể áp dụng.", .65, 2.42, 5.55, .78, 15, MUTED)
for i, (num, label) in enumerate([("RÕ", "Ngôn ngữ gần gũi"), ("GỌN", "Chắt lọc điều quan trọng"), ("HỮU ÍCH", "Gắn với đời sống")]):
    y = 3.45 + i * .86
    rect(s, .65, y, 5.52, .66, WHITE, radius=True, line=LINE)
    text(s, num, .88, y + .19, 1.1, .2, 10, BLUE, True)
    text(s, label, 2.0, y + .16, 3.85, .25, 13, INK, True)
add_image_crop(s, ROOT / "docs/website-homepage-2026.png", 7.05, .58, 5.65, 6.2)
rect(s, 7.05, .58, 5.65, 6.2, WHITE, radius=True, line=LINE, transparency=100)
footer(s, 2)

# Slide 3 — Experience
s = prs.slides.add_slide(blank)
rect(s, 0, 0, 13.333, 7.5, WHITE)
text(s, "02 / TRẢI NGHIỆM", .62, .48, 2.5, .25, 10, BLUE, True)
text(s, "Đọc nhanh. Hiểu sâu. Tìm dễ.", .62, .9, 7.7, .55, 28, INK, True)
text(s, "Một giao diện hiện đại, tập trung vào nội dung.", .64, 1.55, 6.3, .32, 14, MUTED)
add_image_crop(s, ROOT / "tmp/audit-home-desktop.png", .62, 2.12, 8.15, 4.45)
rect(s, .62, 2.12, 8.15, 4.45, WHITE, radius=True, line=LINE, transparency=100)
items = [("TIN NỔI BẬT", "Điểm tin quan trọng ngay đầu trang"), ("CHUYÊN MỤC", "AI · Thiết bị · An ninh · Lập trình"), ("RESPONSIVE", "Mượt mà trên desktop và mobile")]
for i, (head, body) in enumerate(items):
    y = 2.15 + i * 1.43
    rect(s, 9.15, y, 3.55, 1.08, PALE, radius=True, line=LINE)
    rect(s, 9.42, y + .25, .11, .56, BLUE, radius=True)
    text(s, head, 9.72, y + .2, 2.55, .22, 10, BLUE, True)
    text(s, body, 9.72, y + .5, 2.55, .34, 11, INK)
footer(s, 3)

# Slide 4 — Product ecosystem
s = prs.slides.add_slide(blank)
rect(s, 0, 0, 13.333, 7.5, NAVY)
text(s, "03 / HỆ SINH THÁI", .62, .48, 2.7, .25, 10, CYAN, True)
text(s, "Không chỉ là một trang tin.", .62, .9, 6.7, .55, 28, WHITE, True)
text(s, "Nội dung, tài khoản và tương tác trong cùng một trải nghiệm.", .64, 1.55, 7.8, .32, 14, RGBColor(178, 192, 220))
cards = [
    (ROOT / "tmp/audit-article-desktop.png", "BÀI VIẾT CHUYÊN SÂU", "Trình bày rõ ràng, có mục lục và gợi ý liên quan"),
    (ROOT / "public/assets/login-tech-news-hero.webp", "TÀI KHOẢN CÁ NHÂN", "Đăng nhập, lưu bài và tham gia bình luận"),
    (ROOT / "tmp/audit-home-mobile.png", "MOBILE FIRST", "Trải nghiệm nhất quán trên mọi màn hình"),
]
for i, (img, head, body) in enumerate(cards):
    x = .62 + i * 4.17
    add_image_crop(s, img, x, 2.18, 3.78, 2.42)
    rect(s, x, 4.6, 3.78, 1.58, RGBColor(22, 33, 61), radius=True, line=RGBColor(39, 55, 91))
    text(s, head, x + .25, 4.92, 3.2, .23, 10, CYAN, True)
    text(s, body, x + .25, 5.31, 3.2, .52, 11, WHITE)
footer(s, 4, True)

# Slide 5 — Closing
s = prs.slides.add_slide(blank)
add_image_crop(s, ROOT / "public/assets/post-ai-agent.webp", 7.55, 0, 5.783, 7.5)
rect(s, 0, 0, 8.15, 7.5, WHITE)
rect(s, 7.5, 0, .7, 7.5, WHITE, transparency=38)
text(s, "04 / GIÁ TRỊ", .68, .62, 2.3, .25, 10, BLUE, True)
text(s, "Một nền tảng gọn nhẹ,\nsẵn sàng phát triển.", .68, 1.15, 6.1, 1.2, 30, INK, True)
text(s, "Nội dung có chọn lọc", .72, 2.85, 2.55, .3, 14, INK, True)
text(s, "Trải nghiệm thân thiện", 3.52, 2.85, 2.55, .3, 14, INK, True)
rect(s, .72, 3.3, 2.38, .08, BLUE)
rect(s, 3.52, 3.3, 2.38, .08, CYAN)
text(s, "Đọc đúng điều cần biết.", .72, 3.62, 2.45, .4, 12, MUTED)
text(s, "Dùng tốt trên mọi thiết bị.", 3.52, 3.62, 2.65, .4, 12, MUTED)
text(s, "congnghethuongnhat.netlify.app", .72, 5.02, 5.7, .35, 16, BLUE, True)
text(s, "CẢM ƠN", .72, 5.65, 2.2, .28, 11, MUTED, True)
footer(s, 5)

prs.core_properties.title = "Công Nghệ Thường Nhật — Giới thiệu website"
prs.core_properties.subject = "Bài thuyết trình 5 slide, ưu tiên hình ảnh"
prs.core_properties.author = "Công Nghệ Thường Nhật"
prs.save(OUT)
print(OUT.name)
